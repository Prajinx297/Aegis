import os
import shutil
from math import ceil

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = r"D:\Projects\Aegis\deck"
TEMPLATE = r"C:\Users\praji\Downloads\[EXT] Solution Challenge 2026 - Prototype PPT Template.pptx"
OUTPUT = os.path.join(ROOT, "AEGIS_Solution_Challenge_2026_Final.pptx")

PRIMARY = RGBColor(66, 133, 244)
GREEN = RGBColor(52, 168, 83)
YELLOW = RGBColor(251, 188, 5)
RED = RGBColor(234, 67, 53)
INK = RGBColor(24, 34, 52)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(107, 114, 128)
LIGHT = RGBColor(243, 247, 255)
SOFT_BLUE = RGBColor(232, 240, 254)
SOFT_GREEN = RGBColor(232, 245, 233)
SOFT_RED = RGBColor(252, 232, 230)
SOFT_YELLOW = RGBColor(255, 244, 224)
DARK_UI = RGBColor(10, 15, 30)
MID_UI = RGBColor(17, 24, 39)
BORDER = RGBColor(214, 223, 235)


def remove_shape(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)


def clear_slide_content(slide, keep_pictures=True):
    for shape in list(slide.shapes):
        if keep_pictures and shape.shape_type == 13:
            continue
        remove_shape(slide, shape)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=20,
    color=TEXT,
    bold=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    linespacing=1.1,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    p.line_spacing = linespacing
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    font_size=18,
    color=TEXT,
    bullet_color=PRIMARY,
    indent=0.18,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.left_margin = Inches(indent)
        p.indent = Inches(-0.12)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_rule(slide, left, top, width, color=PRIMARY, weight=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_round_rect(slide, left, top, width, height, fill_color, line_color=None, radius=True, line_width=1.2):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    shape.line.width = Pt(line_width)
    return shape


def add_label_chip(slide, left, top, text, fill_color, text_color=TEXT, width=None):
    width = width or Inches(1.6)
    chip = add_round_rect(slide, left, top, width, Inches(0.34), fill_color, fill_color)
    box = chip.text_frame
    box.clear()
    p = box.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    box.vertical_anchor = MSO_ANCHOR.MIDDLE
    return chip


def add_metric(slide, left, top, width, label, value, accent, subtitle):
    add_round_rect(slide, left, top, width, Inches(1.26), RGBColor(255, 255, 255), BORDER)
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), Inches(1.26))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    add_text(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.24), Inches(0.22), label, 10, MUTED, True)
    add_text(slide, left + Inches(0.18), top + Inches(0.33), width - Inches(0.24), Inches(0.45), value, 24, INK, True)
    add_text(slide, left + Inches(0.18), top + Inches(0.82), width - Inches(0.24), Inches(0.28), subtitle, 9.5, MUTED)


def set_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.72), Inches(5.8), Inches(0.42), title, 23, INK, True)
    add_rule(slide, Inches(0.5), Inches(1.17), Inches(1.0), PRIMARY, 2.2)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.28), Inches(7.5), Inches(0.35), subtitle, 11, MUTED)


def add_section_header(slide, kicker, title, body, title_width=Inches(5.2)):
    add_label_chip(slide, Inches(0.55), Inches(0.82), kicker, SOFT_BLUE, PRIMARY, Inches(1.6))
    add_text(slide, Inches(0.55), Inches(1.25), title_width, Inches(0.75), title, 26, INK, True)
    add_text(slide, Inches(0.55), Inches(1.98), Inches(4.8), Inches(0.7), body, 13.5, MUTED)


def draw_cover(slide):
    clear_slide_content(slide, keep_pictures=True)
    add_text(slide, Inches(0.55), Inches(1.0), Inches(4.2), Inches(0.7), "AEGIS", 34, INK, True)
    add_text(
        slide,
        Inches(0.55),
        Inches(1.58),
        Inches(4.9),
        Inches(0.7),
        "Adaptive Embedded Guardian for Integrity in Sports Media",
        16,
        PRIMARY,
        True,
    )
    add_text(
        slide,
        Inches(0.55),
        Inches(2.28),
        Inches(4.75),
        Inches(0.9),
        "An AI-native rights enforcement platform that fingerprints broadcast assets, detects unauthorized redistribution, and generates evidence-backed takedowns in seconds.",
        13,
        TEXT,
    )
    add_label_chip(slide, Inches(0.55), Inches(3.25), "Powered by Gemini 1.5 Pro", SOFT_GREEN, GREEN, Inches(2.2))
    add_label_chip(slide, Inches(2.9), Inches(3.25), "Google Cloud + Vertex AI", SOFT_BLUE, PRIMARY, Inches(2.1))

    ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(6.65), Inches(1.2), Inches(2.25), Inches(2.25))
    ring.fill.background()
    ring.line.color.rgb = PRIMARY
    ring.line.width = Pt(4)
    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(7.15), Inches(1.7), Inches(1.25), Inches(1.25))
    inner.fill.solid()
    inner.fill.fore_color.rgb = SOFT_BLUE
    inner.line.color.rgb = PRIMARY
    inner.line.width = Pt(1.5)
    add_text(slide, Inches(7.24), Inches(2.0), Inches(1.0), Inches(0.45), "AI", 18, PRIMARY, True, align=PP_ALIGN.CENTER)
    for x, y, color in [(6.4, 1.95, PRIMARY), (8.75, 1.95, GREEN), (7.56, 1.0, YELLOW), (7.56, 3.0, RED)]:
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.24), Inches(0.24))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
    add_text(slide, Inches(6.0), Inches(3.7), Inches(2.9), Inches(0.45), "Protect the clip before piracy becomes distribution.", 11.5, MUTED, False, align=PP_ALIGN.CENTER)

    metrics = [
        ("ASSETS REGISTERED", "12", PRIMARY, "fingerprinted and searchable"),
        ("AVG DETECTION TIME", "2.1s", GREEN, "from upload to verdict"),
        ("VALUE PROTECTED", "$2.3M", RED, "estimated rights value preserved"),
    ]
    x_positions = [Inches(0.55), Inches(3.42), Inches(6.29)]
    for x, (label, value, accent, subtitle) in zip(x_positions, metrics):
        add_metric(slide, x, Inches(4.15), Inches(2.55), label, value, accent, subtitle)


def draw_team_slide(slide):
    clear_slide_content(slide, keep_pictures=True)
    add_text(slide, Inches(0.55), Inches(0.9), Inches(2.8), Inches(0.42), "Team Details", 24, INK, True)
    add_rule(slide, Inches(0.55), Inches(1.34), Inches(1.0), PRIMARY, 2.2)

    add_round_rect(slide, Inches(0.55), Inches(1.7), Inches(4.5), Inches(2.65), RGBColor(255, 255, 255), BORDER)
    add_label_chip(slide, Inches(0.78), Inches(1.92), "Submission Snapshot", SOFT_BLUE, PRIMARY, Inches(1.8))
    fields = [
        ("Team name", "AEGIS"),
        ("Team leader", "Submission contact in challenge portal"),
        ("Problem statement", "Broadcasters lose revenue and trust when premium match clips are reposted or deepfaked before legal teams can respond."),
    ]
    y = Inches(2.28)
    for label, value in fields:
        add_text(slide, Inches(0.82), y, Inches(1.25), Inches(0.24), label.upper(), 9.5, MUTED, True)
        add_text(
            slide,
            Inches(2.0),
            y - Inches(0.01),
            Inches(2.65),
            Inches(0.62 if label == "Problem statement" else 0.42),
            value,
            12.2 if label == "Problem statement" else 13,
            TEXT,
            label == "Team name",
        )
        y += Inches(0.55 if label != "Problem statement" else 0.72)

    add_round_rect(slide, Inches(5.35), Inches(1.7), Inches(3.55), Inches(2.65), LIGHT, BORDER)
    add_text(slide, Inches(5.62), Inches(1.95), Inches(2.8), Inches(0.32), "Why this matters", 14, INK, True)
    add_bullets(
        slide,
        Inches(5.55),
        Inches(2.35),
        Inches(2.9),
        Inches(1.55),
        [
            "Live sports value disappears fast once clips spread to mirrors and private channels.",
            "Watermarks alone fail after cropping, reposting, or generative edits.",
            "Broadcasters need proof, not just detection, to enforce at speed.",
        ],
        12.5,
        TEXT,
        PRIMARY,
    )


def draw_solution_brief(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Brief about your solution", "A 3-step workflow designed for judges to understand instantly")
    add_text(slide, Inches(0.6), Inches(1.52), Inches(4.8), Inches(0.42), "Register -> Monitor -> Enforce", 22, INK, True)
    add_text(
        slide,
        Inches(0.6),
        Inches(2.0),
        Inches(4.9),
        Inches(0.72),
        "AEGIS turns every protected asset into a searchable fingerprint, scans suspicious uploads, then uses Gemini 1.5 Pro to confirm violations and draft evidence-backed takedowns.",
        12.5,
        MUTED,
    )
    steps = [
        ("01", "Register", "Extract frames, generate embeddings, index in FAISS, anchor SHA-3 proof on Polygon."),
        ("02", "Monitor", "Scan platform links and suspected uploads, then rank nearest candidates by perceptual similarity."),
        ("03", "Enforce", "Use Gemini to verify ambiguous matches, flag deepfakes, and generate DMCA-ready legal notices."),
    ]
    y = Inches(2.9)
    for idx, (num, label, desc) in enumerate(steps):
        x = Inches(0.6 + idx * 3.0)
        add_round_rect(slide, x, y, Inches(2.6), Inches(1.55), RGBColor(255, 255, 255), BORDER)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.18), y + Inches(0.18), Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = [PRIMARY, GREEN, RED][idx]
        badge.line.fill.background()
        tf = badge.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)
        add_text(slide, x + Inches(0.72), y + Inches(0.16), Inches(1.7), Inches(0.25), label, 14, INK, True)
        add_text(slide, x + Inches(0.18), y + Inches(0.68), Inches(2.18), Inches(0.62), desc, 11.5, MUTED)
    add_text(slide, Inches(6.35), Inches(1.55), Inches(2.1), Inches(0.25), "AI stack", 11, MUTED, True)
    add_label_chip(slide, Inches(6.35), Inches(1.88), "Gemini 1.5 Pro", SOFT_GREEN, GREEN, Inches(1.8))
    add_label_chip(slide, Inches(6.35), Inches(2.3), "ResNet-50 Embeddings", SOFT_BLUE, PRIMARY, Inches(2.0))
    add_label_chip(slide, Inches(6.35), Inches(2.72), "FAISS Vector Search", SOFT_YELLOW, INK, Inches(1.8))
    add_label_chip(slide, Inches(6.35), Inches(3.14), "Polygon Hash Proof", SOFT_RED, RED, Inches(1.8))


def draw_opportunities(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Opportunities", "Where AEGIS is stronger than a basic anti-piracy workflow")
    columns = [
        ("Why different", PRIMARY, "Most tools depend on visible watermarks or simple keyword matching. AEGIS still finds content after cropping, recompression, or reposting."),
        ("Why it solves the problem", GREEN, "FAISS narrows the search fast, Gemini reasons over ambiguous cases, and the system packages proof for action instead of just raising an alert."),
        ("USP", RED, "A courtroom-ready evidence chain: perceptual match, AI reasoning, deepfake assessment, and on-chain integrity proof in one flow."),
    ]
    for i, (title, color, body) in enumerate(columns):
        x = Inches(0.6 + i * 3.05)
        add_round_rect(slide, x, Inches(1.75), Inches(2.75), Inches(2.3), RGBColor(255, 255, 255), BORDER)
        strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(1.75), Inches(0.1), Inches(2.3))
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()
        add_text(slide, x + Inches(0.22), Inches(1.96), Inches(2.2), Inches(0.3), title, 14, INK, True)
        add_text(slide, x + Inches(0.22), Inches(2.36), Inches(2.28), Inches(1.25), body, 12, MUTED)
    add_round_rect(slide, Inches(0.6), Inches(4.25), Inches(8.2), Inches(0.72), LIGHT, BORDER)
    add_text(slide, Inches(0.82), Inches(4.47), Inches(7.8), Inches(0.2), "Why now: generative edits are making sports piracy harder to prove, not just harder to find. Judges need to see that AEGIS closes both gaps.", 12.5, INK, True)


def draw_features(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "List of features offered by the solution", "Lean product surface, high technical depth")
    items = [
        ("Realtime command center", PRIMARY, "Track live violation spikes and severity in one view."),
        ("Perceptual fingerprinting", GREEN, "ResNet-based embeddings survive resize, crop, and re-encode attacks."),
        ("FAISS nearest-neighbor search", YELLOW, "Fast retrieval of likely matches from the registered asset index."),
        ("Gemini verification layer", RED, "Reason over borderline matches and explain the enforcement decision."),
        ("Deepfake anomaly screening", PRIMARY, "Flag manipulated media before it spreads as fake evidence."),
        ("DMCA + audit trail", GREEN, "Generate legal notices with hash proof, timestamps, and case history."),
    ]
    positions = [(0.6, 1.75), (3.45, 1.75), (6.3, 1.75), (0.6, 3.4), (3.45, 3.4), (6.3, 3.4)]
    for (x, y), (title, color, desc) in zip(positions, items):
        x = Inches(x)
        y = Inches(y)
        add_round_rect(slide, x, y, Inches(2.45), Inches(1.35), RGBColor(255, 255, 255), BORDER)
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.18), y + Inches(0.18), Inches(0.34), Inches(0.34))
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()
        add_text(slide, x + Inches(0.62), y + Inches(0.13), Inches(1.6), Inches(0.24), title, 12, INK, True)
        add_text(slide, x + Inches(0.18), y + Inches(0.6), Inches(2.0), Inches(0.48), desc, 10.5, MUTED)


def draw_process_flow(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Process flow diagram", "A judge-friendly view of the full detection pipeline")
    steps = [
        ("Register asset", SOFT_BLUE, PRIMARY),
        ("Extract frames", SOFT_GREEN, GREEN),
        ("Create embeddings", SOFT_YELLOW, INK),
        ("Search FAISS", SOFT_BLUE, PRIMARY),
        ("Verify with Gemini", SOFT_RED, RED),
        ("Generate DMCA", SOFT_GREEN, GREEN),
    ]
    start_x = 0.4
    y = 2.2
    box_w = 1.35
    for i, (label, fill_color, accent) in enumerate(steps):
        x = Inches(start_x + i * 1.5)
        add_round_rect(slide, x, Inches(y), Inches(box_w), Inches(1.1), fill_color, accent)
        add_text(slide, x + Inches(0.12), Inches(y + 0.2), Inches(box_w - 0.24), Inches(0.55), label, 11.5, INK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(box_w), Inches(y + 0.55), x + Inches(1.48), Inches(y + 0.55))
            line.line.color.rgb = MUTED
            line.line.width = Pt(1.5)
    add_round_rect(slide, Inches(0.8), Inches(3.95), Inches(7.8), Inches(0.9), RGBColor(255, 255, 255), BORDER)
    add_text(slide, Inches(1.0), Inches(4.18), Inches(2.3), Inches(0.26), "Operational target", 10.5, MUTED, True)
    add_text(slide, Inches(1.0), Inches(4.42), Inches(1.2), Inches(0.4), "2.1s", 22, PRIMARY, True)
    add_text(slide, Inches(2.2), Inches(4.42), Inches(5.5), Inches(0.3), "average time from suspicious upload ingestion to AI-backed enforcement verdict in the prototype flow", 11.5, TEXT)


def draw_wireframes(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Wireframes / Mock diagrams", "The product stays intuitive by limiting itself to three core actions")
    labels = [("Register", 0.65), ("Monitor", 3.42), ("Enforce", 6.18)]
    for label, xpos in labels:
        x = Inches(xpos)
        add_text(slide, x, Inches(1.55), Inches(2.25), Inches(0.26), label, 13, INK, True)
        frame = add_round_rect(slide, x, Inches(1.85), Inches(2.2), Inches(2.55), RGBColor(250, 250, 251), BORDER)
        frame.line.width = Pt(1.2)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(1.85), Inches(2.2), Inches(0.26))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(245, 247, 250)
        bar.line.fill.background()
        for idx in range(3):
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.18), Inches(2.2 + idx * 0.42), Inches(1.84), Inches(0.16))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(226, 232, 240)
            line.line.fill.background()
        if label == "Register":
            drop = add_round_rect(slide, x + Inches(0.22), Inches(3.1), Inches(1.76), Inches(0.9), RGBColor(245, 247, 250), BORDER)
            drop.line.dash_style = 1
            add_text(slide, x + Inches(0.4), Inches(3.42), Inches(1.3), Inches(0.18), "Upload zone", 10.5, MUTED, True, align=PP_ALIGN.CENTER)
        elif label == "Monitor":
            for row in range(4):
                add_round_rect(slide, x + Inches(0.16), Inches(2.96 + row * 0.28), Inches(1.88), Inches(0.2), RGBColor(243, 247, 255), BORDER)
            add_round_rect(slide, x + Inches(1.55), Inches(2.16), Inches(0.36), Inches(0.36), SOFT_RED, RED)
        else:
            code = add_round_rect(slide, x + Inches(0.18), Inches(2.94), Inches(1.85), Inches(0.78), RGBColor(246, 248, 250), BORDER)
            add_text(slide, x + Inches(0.28), Inches(3.16), Inches(1.6), Inches(0.2), '{"match": true}', 10, MUTED, False, font_name="Consolas")
            add_round_rect(slide, x + Inches(0.18), Inches(3.84), Inches(1.0), Inches(0.32), SOFT_RED, RED)
            add_text(slide, x + Inches(0.28), Inches(3.9), Inches(0.8), Inches(0.16), "DMCA", 10, RED, True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.65), Inches(4.62), Inches(8.0), Inches(0.25), "Each screen is designed to answer one question only: what is protected, what is suspicious, and what action should be taken now.", 11.5, MUTED)


def draw_architecture(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Architecture diagram", "Cloud-native services with a clear AI verification layer")
    add_round_rect(slide, Inches(0.55), Inches(2.05), Inches(1.5), Inches(1.15), SOFT_BLUE, PRIMARY)
    add_text(slide, Inches(0.72), Inches(2.36), Inches(1.15), Inches(0.35), "React dashboard\nBroadcast ops", 13, INK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_round_rect(slide, Inches(2.45), Inches(2.05), Inches(1.55), Inches(1.15), RGBColor(255, 255, 255), BORDER)
    add_text(slide, Inches(2.68), Inches(2.36), Inches(1.1), Inches(0.35), "API + Auth\nCloud Run", 13, INK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_round_rect(slide, Inches(4.35), Inches(1.45), Inches(1.7), Inches(0.8), SOFT_GREEN, GREEN)
    add_text(slide, Inches(4.57), Inches(1.66), Inches(1.25), Inches(0.25), "Gemini 1.5 Pro\nVertex AI", 12, INK, True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, Inches(4.35), Inches(2.45), Inches(1.7), Inches(0.8), SOFT_BLUE, PRIMARY)
    add_text(slide, Inches(4.63), Inches(2.67), Inches(1.1), Inches(0.25), "FAISS\nSearch", 12, INK, True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, Inches(4.35), Inches(3.45), Inches(1.7), Inches(0.8), SOFT_YELLOW, INK)
    add_text(slide, Inches(4.52), Inches(3.67), Inches(1.35), Inches(0.25), "FFmpeg +\nEmbedding Worker", 12, INK, True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, Inches(6.45), Inches(1.2), Inches(1.85), Inches(0.72), RGBColor(255, 255, 255), BORDER)
    add_text(slide, Inches(6.68), Inches(1.42), Inches(1.35), Inches(0.2), "Cloud Storage", 12, INK, True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, Inches(6.45), Inches(2.22), Inches(1.85), Inches(0.72), RGBColor(255, 255, 255), BORDER)
    add_text(slide, Inches(6.76), Inches(2.44), Inches(1.22), Inches(0.2), "Firestore", 12, INK, True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, Inches(6.45), Inches(3.24), Inches(1.85), Inches(0.72), RGBColor(255, 255, 255), BORDER)
    add_text(slide, Inches(6.6), Inches(3.46), Inches(1.5), Inches(0.2), "Polygon proof layer", 12, INK, True, align=PP_ALIGN.CENTER)

    for x1, y1, x2, y2 in [
        (2.05, 2.62, 2.45, 2.62),
        (4.0, 2.62, 4.35, 2.62),
        (4.0, 2.62, 4.35, 1.86),
        (4.0, 2.62, 4.35, 3.86),
        (6.05, 1.86, 6.45, 1.56),
        (6.05, 2.62, 6.45, 2.58),
        (6.05, 3.86, 6.45, 3.6),
    ]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = MUTED
        line.line.width = Pt(1.2)

    add_text(slide, Inches(0.7), Inches(4.4), Inches(7.8), Inches(0.26), "Platform connectors ingest suspected URLs from public social channels and mirror sites; verified cases are written back into the audit trail.", 11.5, MUTED)


def draw_tech_slide(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Technologies to be used in the solution", "Every technology maps directly to a part of the problem")
    groups = [
        ("Frontend", PRIMARY, ["React", "Tailwind CSS", "Firebase Hosting"]),
        ("AI / ML", GREEN, ["Gemini 1.5 Pro", "Vertex AI", "ResNet-50", "FAISS"]),
        ("Security", RED, ["Firebase Auth", "SHA-3 hashing", "Polygon"]),
        ("Infra", YELLOW, ["Cloud Run", "Cloud Storage", "Firestore", "FFmpeg"]),
    ]
    for i, (name, accent, techs) in enumerate(groups):
        x = Inches(0.75 + i * 2.2)
        add_round_rect(slide, x, Inches(1.9), Inches(1.95), Inches(2.4), RGBColor(255, 255, 255), BORDER)
        add_text(slide, x + Inches(0.18), Inches(2.12), Inches(1.45), Inches(0.22), name, 13, accent, True)
        y = Inches(2.55)
        for tech in techs:
            add_label_chip(slide, x + Inches(0.18), y, tech, LIGHT, TEXT, Inches(1.55))
            y += Inches(0.42)


def draw_cost_slide(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Estimated implementation cost", "Prototype-friendly today, scalable for league deployment later")
    add_round_rect(slide, Inches(0.75), Inches(1.85), Inches(4.45), Inches(2.7), RGBColor(255, 255, 255), BORDER)
    add_text(slide, Inches(0.98), Inches(2.05), Inches(1.8), Inches(0.22), "Monthly estimate", 13, INK, True)
    costs = [
        ("Cloud Run API + workers", "$70"),
        ("Storage + network egress", "$45"),
        ("Gemini verification calls", "$180"),
        ("FAISS + background monitoring", "$95"),
        ("Logging / alerting", "$40"),
    ]
    y = Inches(2.45)
    for label, value in costs:
        add_text(slide, Inches(1.0), y, Inches(3.0), Inches(0.2), label, 11.5, TEXT)
        add_text(slide, Inches(4.15), y, Inches(0.6), Inches(0.2), value, 11.5, INK, True, align=PP_ALIGN.RIGHT)
        add_rule(slide, Inches(1.0), y + Inches(0.23), Inches(3.75), BORDER, 0.8)
        y += Inches(0.37)
    add_text(slide, Inches(1.0), Inches(4.02), Inches(3.0), Inches(0.24), "Estimated prototype total", 12, MUTED, True)
    add_text(slide, Inches(3.75), Inches(3.98), Inches(1.0), Inches(0.22), "$430/mo", 13.5, PRIMARY, True, align=PP_ALIGN.RIGHT)
    add_round_rect(slide, Inches(5.55), Inches(1.85), Inches(3.05), Inches(2.7), SOFT_BLUE, PRIMARY)
    add_text(slide, Inches(5.82), Inches(2.08), Inches(2.4), Inches(0.25), "Why the cost stays efficient", 13, INK, True)
    add_bullets(
        slide,
        Inches(5.75),
        Inches(2.45),
        Inches(2.45),
        Inches(1.6),
        [
            "Most assets are resolved by fast vector search before Gemini is invoked.",
            "Event-driven workers keep heavy compute off until a new asset or suspected violation arrives.",
            "Cloud services scale with match volume rather than requiring fixed hardware.",
        ],
        11.5,
        TEXT,
        PRIMARY,
    )


def draw_mvp_snapshots(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Snapshots of the MVP", "Two key product moments from the interactive prototype")
    for i, title in enumerate(["Command Center", "Detect Violation"]):
        x = Inches(0.62 + i * 4.28)
        screen = add_round_rect(slide, x, Inches(1.75), Inches(3.75), Inches(2.72), DARK_UI, MID_UI)
        screen.line.width = Pt(1.3)
        top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(1.75), Inches(3.75), Inches(0.28))
        top.fill.solid()
        top.fill.fore_color.rgb = MID_UI
        top.line.fill.background()
        for dot_idx, color in enumerate([RED, YELLOW, GREEN]):
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.16 + dot_idx * 0.16), Inches(1.84), Inches(0.08), Inches(0.08))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
        add_text(slide, x + Inches(0.32), Inches(2.12), Inches(2.9), Inches(0.24), title, 12.5, RGBColor(255, 255, 255), True)
        if i == 0:
            for c in range(3):
                add_round_rect(slide, x + Inches(0.2 + c * 1.15), Inches(2.45), Inches(0.95), Inches(0.52), MID_UI, MID_UI)
            for row in range(5):
                add_round_rect(slide, x + Inches(0.2), Inches(3.14 + row * 0.26), Inches(2.35), Inches(0.18), RGBColor(24, 34, 52), RGBColor(24, 34, 52))
            add_round_rect(slide, x + Inches(2.72), Inches(2.45), Inches(0.68), Inches(1.96), MID_UI, MID_UI)
        else:
            add_round_rect(slide, x + Inches(0.22), Inches(2.48), Inches(1.3), Inches(0.58), MID_UI, MID_UI)
            add_round_rect(slide, x + Inches(1.7), Inches(2.48), Inches(1.68), Inches(0.58), MID_UI, MID_UI)
            meter_bg = add_round_rect(slide, x + Inches(0.22), Inches(3.18), Inches(3.08), Inches(0.18), RGBColor(24, 34, 52), RGBColor(24, 34, 52))
            meter = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Inches(0.22), Inches(3.18), Inches(2.55), Inches(0.18))
            meter.fill.solid()
            meter.fill.fore_color.rgb = RED
            meter.line.fill.background()
            add_round_rect(slide, x + Inches(0.22), Inches(3.55), Inches(3.08), Inches(0.72), MID_UI, MID_UI)
            add_text(slide, x + Inches(0.35), Inches(3.8), Inches(2.5), Inches(0.2), '{"match": true, "confidence": 0.94}', 9.5, RGBColor(226, 232, 240), False, font_name="Consolas")
        add_text(slide, x + Inches(0.22), Inches(4.6), Inches(3.2), Inches(0.2), "Minimal UI, fast action paths, and visible AI evidence.", 10.5, RGBColor(200, 208, 219))


def draw_future_slide(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Additional details / Future development", "A narrow MVP today, a league-scale product path tomorrow")
    milestones = [
        ("Now", PRIMARY, "Upload registration, FAISS search, Gemini verification, DMCA drafting"),
        ("Next", GREEN, "Automated crawlers, multilingual legal packs, human reviewer feedback loop"),
        ("Later", RED, "Live-stream fingerprinting, broadcaster APIs, rights marketplace integrations"),
    ]
    for i, (label, accent, body) in enumerate(milestones):
        x = Inches(0.75 + i * 2.8)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.48), Inches(2.2), Inches(0.48), Inches(0.48))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        add_text(slide, x, Inches(2.8), Inches(1.9), Inches(0.24), label, 13, INK, True, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(3.18), Inches(1.9), Inches(0.8), body, 11.5, MUTED, align=PP_ALIGN.CENTER)
        if i < 2:
            add_rule(slide, x + Inches(1.9), Inches(2.44), Inches(0.8), BORDER, 1.2)


def draw_links_slide(slide):
    clear_slide_content(slide, keep_pictures=True)
    set_title(slide, "Submission links", "Slide is ready for the final URLs once publishing is complete")
    rows = [
        ("GitHub public repository", "Add final public repo URL before submission"),
        ("Demo video link (3 minutes)", "Add uploaded demo URL"),
        ("MVP link", "Attach deployed prototype URL"),
        ("Working prototype link", "Attach production or staging URL"),
    ]
    y = Inches(1.8)
    for label, value in rows:
        add_round_rect(slide, Inches(0.78), y, Inches(7.95), Inches(0.72), RGBColor(255, 255, 255), BORDER)
        add_text(slide, Inches(1.02), y + Inches(0.18), Inches(2.45), Inches(0.2), label, 12, INK, True)
        add_text(slide, Inches(3.7), y + Inches(0.18), Inches(4.5), Inches(0.2), value, 11.5, MUTED)
        y += Inches(0.88)
    add_text(slide, Inches(0.82), Inches(5.1), Inches(7.4), Inches(0.22), "If you send the final links, this slide can be patched in a minute without changing the rest of the deck.", 10.5, MUTED)


def draw_impact_slide(slide):
    clear_slide_content(slide, keep_pictures=True)
    add_text(slide, Inches(0.65), Inches(1.0), Inches(3.8), Inches(0.3), "EXPECTED IMPACT", 12, PRIMARY, True)
    add_text(slide, Inches(0.65), Inches(1.34), Inches(4.25), Inches(1.15), "78% faster review-to-action for sports media enforcement", 23, INK, True)
    add_text(slide, Inches(0.65), Inches(2.55), Inches(3.9), Inches(0.46), "AEGIS reduces manual triage and packages proof before clips spread across mirrors.", 12.5, MUTED)

    chart_data = ChartData()
    chart_data.categories = ["Manual review", "AEGIS workflow"]
    chart_data.add_series("Minutes", (9.5, 2.1))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(5.2), Inches(1.45), Inches(3.2), Inches(2.2), chart_data).chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 10
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 2
    chart.value_axis.format.line.color.rgb = BORDER
    chart.category_axis.format.line.color.rgb = BORDER
    chart.chart_title.has_text_frame = False
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = PRIMARY

    add_metric(slide, Inches(0.7), Inches(3.7), Inches(2.45), "VIOLATIONS / DAY", "47", RED, "prototype monitoring baseline")
    add_metric(slide, Inches(3.3), Inches(3.7), Inches(2.45), "DEEPFAKES FLAGGED", "9", GREEN, "high-risk cases surfaced")
    add_metric(slide, Inches(5.9), Inches(3.7), Inches(2.45), "EST. VALUE PROTECTED", "$2.3M", PRIMARY, "rights value preserved")


def draw_thank_you(slide):
    # Keep template artwork and add a small project signature.
    add_text(slide, Inches(5.85), Inches(1.95), Inches(2.2), Inches(0.28), "AEGIS", 18, PRIMARY, True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.52), Inches(2.28), Inches(2.8), Inches(0.34), "AI-native sports media integrity", 12, MUTED, align=PP_ALIGN.CENTER)


def rebuild_deck():
    os.makedirs(ROOT, exist_ok=True)
    shutil.copyfile(TEMPLATE, OUTPUT)
    prs = Presentation(OUTPUT)

    draw_cover(prs.slides[0])
    draw_team_slide(prs.slides[1])
    draw_solution_brief(prs.slides[2])
    draw_opportunities(prs.slides[3])
    draw_features(prs.slides[4])
    draw_process_flow(prs.slides[5])
    draw_wireframes(prs.slides[6])
    draw_architecture(prs.slides[7])
    draw_tech_slide(prs.slides[8])
    draw_cost_slide(prs.slides[9])
    draw_mvp_snapshots(prs.slides[10])
    draw_future_slide(prs.slides[11])
    draw_links_slide(prs.slides[12])
    draw_impact_slide(prs.slides[13])
    draw_thank_you(prs.slides[14])

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(rebuild_deck())
